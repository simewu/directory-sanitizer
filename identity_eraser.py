from docx import Document
from PyPDF2 import PdfFileReader, PdfFileWriter
from PyPDF2.generic import NameObject, DecodedStreamObject, EncodedStreamObject
import openpyxl
import os
import pptx
import pypdftk
import re
import shutil
import sys
import tarfile
import tempfile
import zipfile
# Requires installing pdftk (https://www.pdflabs.com/tools/pdftk-the-pdf-toolkit/)

firstName = input('Enter your first name: ').strip().capitalize()
lastName = input('Enter your first name: ').strip().capitalize()
nickname = firstName[0] + lastName
firstNameTo = 'John'
lastNameTo = 'Smith'
nicknameTo = firstNameTo[0] + lastNameTo

# File size threshold check applies to .pdf, .xlsx, and misc files
FileSizeThresholdMegabytes = 30
FileSkipOption = 0 # 0 to ask the user each time, 1 to skip, 2 to process

pairs = []
pairs.append([nickname, nicknameTo])
pairs.append([firstName, firstNameTo])
pairs.append([lastName, lastNameTo])

print()
print(f'Hello, {firstName} {lastName}.')

def replace(string, morePairs=[]):
	string2 = string
	for a, b in morePairs:
		string2 = re.sub(a, b, string2)
		string2 = re.sub(a.upper(), b.upper(), string2)
		string2 = re.sub(a.lower(), b.lower(), string2, flags=re.IGNORECASE)
	for a, b in pairs:
		string2 = re.sub(a, b, string2)
		string2 = re.sub(a.upper(), b.upper(), string2)
		string2 = re.sub(a.lower(), b.lower(), string2, flags=re.IGNORECASE)
	return string2


# Given a regular expression, list the directories that match it, and ask for user input
def selectDir(regex, subdirs = False):
	dirs = []
	if subdirs:
		for (dirpath, dirnames, filenames) in os.walk('.'):
			if dirpath[:2] == '.\\': dirpath = dirpath[2:]
			if bool(re.match(regex, dirpath)):
				dirs.append(dirpath)
	else:
		for obj in os.listdir(os.curdir):
			if os.path.isdir(obj) and bool(re.match(regex, obj)):
				dirs.append(obj)

	print()
	if len(dirs) == 0:
		print(f'No directories were found that match "{regex}"')
		print()
		return ''

	print('List of directories:')
	for i, directory in enumerate(dirs):
		print(f'  Directory {i + 1}  -  {directory}')
	print()

	selection = None
	while selection is None:
		try:
			i = int(input(f'Please select a directory (1 to {len(dirs)}): '))
		except KeyboardInterrupt:
			sys.exit()
		except:
			pass
		if i > 0 and i <= len(dirs):
			selection = dirs[i - 1]
	print()
	return selection

# List the files with a regular expression
def listSubdirs(directory = ''):
	output = []
	for path, subdirs, files in os.walk(directory):
		output.append(path)
	return output

# List the files with a regular expression
def listFiles(directory = ''):
	output = []
	for path, subdirs, files in os.walk(directory):
		for name in files:
			output.append(os.path.join(path, name))
	return output

# Rename occurences within a path
def processPath(path):
	name = os.path.basename(path)
	pathBefore = path[:-len(name) - 1]
	name2 = replace(name)
	if name != name2:
		print(f'Renaming "{path}"...')
		os.rename(path, os.path.join(pathBefore, name2))
		return True
	return False

def sizeTooBigCheck(path):
	global FileSizeThresholdMegabytes, FileSkipOption
	megabytes = os.path.getsize(path) / 1000000
	if megabytes > FileSizeThresholdMegabytes and FileSkipOption == 0:
		print()
		choice = input(f'The file "{path}" is {megabytes} megabytes. This may take a while to process. (s=skip, p=process, s!=skip and don\'t ask again, p!=process and don\'t ask again): ').strip().lower()
		while choice not in ['s', 'p', 's!', 'p!']:
			choice = input(f'(s=skip, p=process, s!=skip and don\'t ask again, p!=process and don\'t ask again): ').strip().lower()
		if choice == 's':
			return False
		if choice == 's!':
			FileSkipOption = 1
			return False
		if choice == 'p':
			return True
		if choice == 'p!':
			FileSkipOption = 2
			return True
	elif FileSkipOption == 1: return False
	return True

# Rename occurences within a PDF
def processPDF(path):
	def PDF_uncompress(path):
		temp_dir = tempfile.TemporaryDirectory()
		temp_path = os.path.join(temp_dir.name, os.path.basename(path))
		pypdftk.uncompress('"' + path + '"', '"' + temp_path + '"')
		shutil.move(temp_path, path)
		temp_dir.cleanup()

	def PDF_compress(path):
		temp_dir = tempfile.TemporaryDirectory()
		temp_path = os.path.join(temp_dir.name, os.path.basename(path))
		pypdftk.compress('"' + path + '"', '"' + temp_path + '"')
		shutil.move(temp_path, path)
		temp_dir.cleanup()

	def PDF_replaceText(content):
		# Created a set of regex/replacements that handle the PDF syntax
		# e.g. "Test" may be encoded as "(T) 73 (e)-20.42(st)"
		_regex1 = ''
		for i in range(len(firstName)):
			if i > 0: _regex1 += '(\) *[-+]?[0-9\.]+ *\()?'
			_regex1 += firstName[i]
		regex1_ = ''
		for i in range(len(firstNameTo)):
			if i > 0: regex1_ += f'\\{i}'
			regex1_ += firstNameTo[i]

		_regex2 = ''
		for i in range(len(lastName)):
			if i > 0: _regex2 += '(\) *[-+]?[0-9\.]+ *\()?'
			_regex2 += lastName[i]
		regex2_ = ''
		for i in range(len(lastNameTo)):
			if i > 0: regex2_ += f'\\{i}'
			regex2_ += lastNameTo[i]

		morePairs = []
		morePairs.append([rf'{_regex1}', rf'{regex1_}'])
		morePairs.append([rf'{_regex2}', rf'{regex2_}'])
		changed = False
		lines = content.splitlines()
		result = ''
		in_text = False
		for line in lines:
			if line == 'BT': in_text = True
			elif line == 'ET': in_text = False
			elif in_text:
				cmd = line[-2:]
				if cmd.lower() == 'tj':
					replaced_line = replace(line, morePairs)
					if replaced_line != line:
						#print('From\t', line)
						#print('To  \t', replaced_line)
						changed = True
					result += replaced_line + '\n'
				else:
					result += line + '\n'
				continue
			result += line + '\n'
		return result, changed

	def PDF_processData(object):
		data = object.getData()
		decoded_data = data.decode('utf-8', errors='ignore')
		replaced_data, changed = PDF_replaceText(decoded_data)
		encoded_data = replaced_data.encode('utf-8')
		if object.decodedSelf is not None:
			object.decodedSelf.setData(encoded_data)
		else:
			object.setData(encoded_data)
		return changed

	try:
		PDF_uncompress(path)
	except: pass
	# Attempt 1
	changed = False
	pdf = PdfFileReader(open(path, 'rb'))
	writer = PdfFileWriter() 
	for page in pdf.pages:
		contents = page.getContents().getData()
		contentsBefore = contents
		for a, b in pairs:
			contents = contents.replace(a.encode('utf-8'), b.encode('utf-8'))
			contents = contents.replace(a.upper().encode('utf-8'), b.upper().encode('utf-8'))
			contents = contents.replace(a.lower().encode('utf-8'), b.lower().encode('utf-8'))
		if contents != contentsBefore:
			changed = True
			page.getContents().setData(contents)
		writer.addPage(page)
	if changed:
		print(f'Attempt 1: Rewriting "{path}"...')
		with open(path, 'wb') as file:
			writer.write(file)
	try:
		changed = False
		pdf = PdfFileReader(path)
		writer = PdfFileWriter()
		for page_number in range(0, pdf.getNumPages()):
			page = pdf.getPage(page_number)
			contents = page.getContents()
			if isinstance(contents, DecodedStreamObject) or isinstance(contents, EncodedStreamObject):
				status = PDF_processData(contents)
				if status: changed = True
			elif len(contents) > 0:
				for obj in contents:
					if isinstance(obj, DecodedStreamObject) or isinstance(obj, EncodedStreamObject):
						streamObj = obj.getObject()
						status = PDF_processData(streamObj)
						if status: changed = True
			try:
				page[NameObject('/Contents')] = contents.decodedSelf
			except: pass
			writer.addPage(page)
		if changed:
			print(f'Attempt 2: Rewriting "{path}"...')
			with open(path, 'wb') as file:
				writer.write(file)
			PDF_compress(path)

	except: pass # PDF is corrupt or encrypted, skip

# Rename occurences within a DOCX
def processDOCX(path):
	def docx_replace_regex(doc_obj, regex, replace, caseSensitive):
		for p in doc_obj.paragraphs:
			if re.search(regex, p.text):
				inline = p.runs
				# Loop added to work with runs (strings with same style)
				for i in range(len(inline)):
					if re.search(regex, inline[i].text):
						if caseSensitive:
							text = re.sub(regex, replace, inline[i].text)
						else:
							text = re.sub(regex, replace, inline[i].text, flags=re.IGNORECASE)
						inline[i].text = text
		for table in doc_obj.tables:
			for row in table.rows:
				for cell in row.cells:
					docx_replace_regex(cell, regex, replace, caseSensitive)
		if 'sections' in dir(doc_obj):
			for section in doc_obj.sections:
				header = section.header
				docx_replace_regex(header, regex, replace, caseSensitive)

	try:
		doc = Document(path)
		for a, b in pairs:
			docx_replace_regex(doc, a, b, True)
			docx_replace_regex(doc, a.upper(), b.upper(), True)
			docx_replace_regex(doc, a.lower(), b.lower(), False)
		temp_dir = tempfile.TemporaryDirectory()
		temp_path = os.path.join(temp_dir.name, os.path.basename(path))
		doc.save(temp_path)
		shutil.move(temp_path, path)
		temp_dir.cleanup()
	except: pass

# Rename occurences within a XLSX
def processXLSX(path):
	workbook = openpyxl.load_workbook(path)
	for worksheet in workbook.worksheets:
		for row in range(1, worksheet.max_row + 1):
			for column in range(1,worksheet.max_column + 1):
				value = worksheet.cell(row,column).value
				if value != None: 
					try:
						worksheet.cell(row,column).value = replace(value)
					except: pass

	temp_dir = tempfile.TemporaryDirectory()
	temp_path = os.path.join(temp_dir.name, os.path.basename(path))
	workbook.save(temp_path)
	shutil.move(temp_path, path)
	temp_dir.cleanup()

# Rename occurences within a PPTX
def processPPTX(path):
	ppt = pptx.Presentation(path)
	for slide in ppt.slides:
		for shape in slide.shapes:
			if shape.has_text_frame:
				for paragraph in shape.text_frame.paragraphs:
					for run in paragraph.runs:
						try:
							run.text = replace(run.text)
						except: pass
	temp_dir = tempfile.TemporaryDirectory()
	temp_path = os.path.join(temp_dir.name, os.path.basename(path))
	ppt.save(temp_path)
	shutil.move(temp_path, path)
	temp_dir.cleanup()

# Rename occurences within a ZIP (recursively unzips, processes, rezips)
def processZIP(path):
	temp_dir = tempfile.TemporaryDirectory()
	try:
		with zipfile.ZipFile(path, 'r') as zip_ref:
			zip_ref.extractall(temp_dir.name)
		try:
			processDirectory(temp_dir.name)
		except: pass
		shutil.make_archive(path[:-4], 'zip', temp_dir.name)
	except: pass
	temp_dir.cleanup()

# Rename occurences within a ZIP (recursively unzips, processes, rezips)
def processTAR(path):
	temp_dir = tempfile.TemporaryDirectory()
	try:
		with tarfile.open(path, 'r') as tar_ref:
			tar_ref.extractall(temp_dir.name)
		try:
			processDirectory(temp_dir.name)
		except: pass
		os.remove(path)
		with tarfile.open(path, 'w:gz') as tar:
			tar.add(temp_dir.name, arcname=os.path.basename(path))
	except: pass
	temp_dir.cleanup()
	
# Processes a directory, looking for each supported file type
def processDirectory(baseDirectory):
	print('[ Renaming directories ]')
	directories = listSubdirs(baseDirectory)
	i = 0
	while i < len(directories):
		directory = directories[i]
		changed = processPath(directory)
		if changed:
			directories = listSubdirs(baseDirectory)
		i += 1

	print('[ Renaming files ]')
	files = listFiles(baseDirectory)
	for path in files:
		processPath(path)

	print('[ Editing direct ASCII files ]')
	files = listFiles(baseDirectory)
	for path in files:
		_, extension = os.path.splitext(path)
		if extension in ['.pdf', '.docx', '.xlsx', '.pptx', '.zip']: continue
		if sizeTooBigCheck(path):
			with open(path, 'r', encoding='utf-8', errors='ignore') as file:
				contents = file.read()
			contents2 = replace(contents)
			if contents != contents2:
				print(f'Editing "{path}"...')
				with open(path, 'w', encoding='utf-8', errors='ignore') as file:
					file.write(contents2)

	print('[ Editing PDF files ]')
	files = listFiles(baseDirectory)
	pdfs = []
	for path in files:
		if path.lower().endswith('.pdf'):
			pdfs.append(path)
	for path in pdfs:
		if sizeTooBigCheck(path):
			processPDF(path)
		else: continue
	
	print('[ Editing DOCX files ]')
	files = listFiles(baseDirectory)
	pdfs = []
	for path in files:
		if path.lower().endswith('.docx'):
			pdfs.append(path)
	for path in pdfs:
		print(f'Editing {path}...')
		processDOCX(path)

	print('[ Editing XLSX files ]')
	files = listFiles(baseDirectory)
	pdfs = []
	for path in files:
		if path.lower().endswith('.xlsx'):
			pdfs.append(path)
	for path in pdfs:
		if sizeTooBigCheck(path):
			print(f'Editing {path}...')
			processXLSX(path)
		else: continue

	print('[ Editing PPTX files ]')
	files = listFiles(baseDirectory)
	pdfs = []
	for path in files:
		if path.lower().endswith('.pptx'):
			pdfs.append(path)
	for path in pdfs:
		print(f'Editing {path}...')
		processPPTX(path)

	print('[ Editing ZIP files ]')
	files = listFiles(baseDirectory)
	pdfs = []
	for path in files:
		if path.lower().endswith('.zip'):
			pdfs.append(path)
	for path in pdfs:
		print(f'Unzipping {path}...')
		processZIP(path)

	print('[ Editing TAR files ]')
	files = listFiles(baseDirectory)
	pdfs = []
	for path in files:
		if path.lower().endswith('.tar'):
			pdfs.append(path)
		if path.lower().endswith('.tar.gz'):
			pdfs.append(path)
		if path.lower().endswith('.tar.bz2'):
			pdfs.append(path)
		if path.lower().endswith('.tar.Z'):
			pdfs.append(path)
		if path.lower().endswith('.tar.xz'):
			pdfs.append(path)
	for path in pdfs:
		print(f'Uncompressing {path}...')
		processTAR(path)

if __name__ == '__main__':
	baseDirectory = selectDir(r'.*', False)
	if not baseDirectory.endswith('_RemovedIdentity'):
		baseDirectoryTo = baseDirectory + '_RemovedIdentity'
		if os.path.exists(baseDirectoryTo):
			print(f'Removing old to "{baseDirectoryTo}"...')
			shutil.rmtree(baseDirectoryTo)
		print(f'Copying to "{baseDirectoryTo}"...')
		shutil.copytree(baseDirectory, baseDirectoryTo)
		baseDirectory = baseDirectoryTo

	processDirectory(baseDirectory)
	print(f'Successfully created "{baseDirectory}".')